#!/usr/bin/env python3
"""
Generate a PPTX summary for the ERT app using existing docs/*.md and docs/*.png
Produces: docs/ERT_presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
import os

WORKDIR = os.path.dirname(os.path.dirname(__file__))
DOCS = os.path.join(WORKDIR, 'docs')
OUT = os.path.join(DOCS, 'ERT_presentation.pptx')

# Files to include (order matters)
FILES = [
    'ert-presentation-slide.md',
    'ert-overview.md',
    'onlot-dataflow.md',
    'onprod-dataflow.md',
    'onscribe-dataflow.md',
    'onslice-dataflow.md',
    'ert-databricks-migration-presentation.md'
]

IMAGES = [
    'ert-overview.png',
    'onlot-dataflow.png',
    'onprod-dataflow.png',
    'onscribe-dataflow.png',
    'onslice-dataflow.png',
    'onwmap-level2-config.png',  # immediately after Slice
    'onwmap-level2-product.png'
]

prs = Presentation()
# Title slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Exensio Reference Tables (ERT) — Review"
subtitle.text = "App overview, flows, and Databricks migration plan"

# Helper to add bullet slide
def add_bullet_slide(title_text, bullets):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1].text_frame
    title.text = title_text
    body.clear()
    for i, b in enumerate(bullets):
        p = body.add_paragraph() if i>0 else body.paragraphs[0]
        p.text = b
        p.level = 0
        p.font.size = Pt(18)
    return slide

# Read MD and extract top header and first paragraph bullets until '---'
def md_to_blocks(path):
    if not os.path.exists(path):
        return None, []
    with open(path,'r',encoding='utf-8') as f:
        lines = [l.rstrip() for l in f]
    # find first header
    header = None
    bullets = []
    hr = False
    for l in lines:
        if not l:
            continue
        if l.startswith('#') and header is None:
            header = l.lstrip('# ').strip()
            continue
        if l.strip() == '---':
            break
        # collect short lines as bullets
        if len(bullets) < 10 and not l.startswith('```'):
            # remove markdown list markers
            t = l.lstrip('- ').lstrip('* ').strip()
            if t:
                bullets.append(t)
    return header or os.path.basename(path), bullets

# Add agenda slide
agenda = ["Title & Context", "Current-state snapshot", "Key flows: OnLot / OnProd / OnScribe / OnSlice / OnWmap", "Integration landscape", "Databricks migration plan", "Next steps"]
add_bullet_slide("Agenda", agenda)

# Add slides from files
for fn in FILES:
    path = os.path.join(DOCS, fn)
    header, bullets = md_to_blocks(path)
    if header is None:
        continue
    # limit bullet length
    bullets = bullets[:6] if bullets else ["See docs/{}".format(fn)]
    add_bullet_slide(header, bullets)

# Add architecture image slide if exists
def add_image_slide(title_text, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = title_text
    left = Inches(0.5)
    top = Inches(1.2)
    if os.path.exists(image_path):
        try:
            slide.shapes.add_picture(image_path, left, top, width=Inches(9))
        except Exception as e:
            # Fallback when image file is corrupt or unsupported by PIL
            tx = slide.shapes.add_textbox(left, top, Inches(9), Inches(1))
            tf = tx.text_frame
            tf.text = "[Could not render image: {}]".format(os.path.basename(image_path))
    else:
        tx = slide.shapes.add_textbox(left, top, Inches(9), Inches(1))
        tf = tx.text_frame
        tf.text = "[Image not found: {}]".format(os.path.basename(image_path))

for img in IMAGES:
    img_path = os.path.join(DOCS, img)
    if os.path.exists(img_path):
        add_image_slide(os.path.splitext(img)[0].replace('-', ' ').upper(), img_path)

# Final: Key migration slide from databricks doc
hdr, bullets = md_to_blocks(os.path.join(DOCS, 'ert-databricks-migration-presentation.md'))
if hdr:
    add_bullet_slide('Databricks Migration: Summary', bullets[:6])

# Save presentation
prs.save(OUT)
print('Wrote', OUT)
