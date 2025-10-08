#!/usr/bin/env python3
"""
Generate a PPTX summary for the ERT app (ExensioRefTables only) using existing docs/*.md and docs/*.png
Produces: docs/ERT_presentation_ert_only.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
import os

WORKDIR = os.path.dirname(os.path.dirname(__file__))
DOCS = os.path.join(WORKDIR, 'docs')
OUT = os.path.join(DOCS, 'ERT_presentation_ert_only.pptx')

FILES = [
    'ert-presentation-slide.md',
    'ert-overview.md',
    'onlot-dataflow.md',
    'onprod-dataflow.md',
    'onscribe-dataflow.md',
    'onslice-dataflow.md'
]

IMAGES = [
    'ert-overview.png',
    'onlot-dataflow.png',
    'onprod-dataflow.png',
    'onscribe-dataflow.png',
    'onslice-dataflow.png',
    'onwmap-level2-config.png',  # place WMAP after Slice
    'onwmap-level2-product.png'
]

prs = Presentation()
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Exensio Reference Tables (ERT) — Review"
slide.placeholders[1].text = "Focused: ExensioRefTables flows and operations"

# helper
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

def add_bullet_slide(title_text, bullets):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1].text_frame
    title.text = title_text
    body.clear()
    for i, b in enumerate(bullets):
        p = body.add_paragraph() if i>0 else body.paragraphs[0]
        p.text = b
        p.level = 0
        p.font.size = Pt(18)
    return slide

# md reader
def md_to_blocks(path):
    if not os.path.exists(path):
        return None, []
    with open(path,'r',encoding='utf-8') as f:
        lines = [l.rstrip() for l in f]
    header = None
    bullets = []
    for l in lines:
        if not l:
            continue
        if l.startswith('#') and header is None:
            header = l.lstrip('# ').strip()
            continue
        if l.strip() == '---':
            break
        t = l.lstrip('- ').lstrip('* ').strip()
        if t and len(bullets) < 8:
            bullets.append(t)
    return header or os.path.basename(path), bullets

agenda = ["Title & Context", "Key flows: OnLot / OnProd / OnScribe / OnSlice / OnWmap", "Integration overview", "Operational notes"]
add_bullet_slide('Agenda', agenda)

for fn in FILES:
    header, bullets = md_to_blocks(os.path.join(DOCS, fn))
    if header:
        add_bullet_slide(header, bullets[:6])

# images
for img in IMAGES:
    img_path = os.path.join(DOCS, img)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = os.path.splitext(img)[0].replace('-', ' ').upper()
    left = Inches(0.5)
    top = Inches(1.2)
    if os.path.exists(img_path):
        try:
            slide.shapes.add_picture(img_path, left, top, width=Inches(9))
        except Exception:
            tx = slide.shapes.add_textbox(left, top, Inches(9), Inches(1))
            tx.text = f"[Could not render image: {img}]"
    else:
        tx = slide.shapes.add_textbox(left, top, Inches(9), Inches(1))
        tx.text = f"[Image not found: {img}]"

prs.save(OUT)
print('Wrote', OUT)
