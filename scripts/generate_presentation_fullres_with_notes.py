#!/usr/bin/env python3
"""
Generate a PPTX that embeds images at full resolution and injects speaker notes from markdown files.
Produces: docs/ERT_presentation_fullres_with_notes.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
import os

WORKDIR = os.path.dirname(os.path.dirname(__file__))
DOCS = os.path.join(WORKDIR, 'docs')
OUT = os.path.join(DOCS, 'ERT_presentation_fullres_with_notes.pptx')

FILES = [
    ('ERT Summary', 'ert-presentation-slide.md'),
    ('Overview', 'ert-overview.md'),
    ('OnLot', 'onlot-dataflow.md'),
    ('OnProd', 'onprod-dataflow.md'),
    ('OnScribe', 'onscribe-dataflow.md'),
    ('OnSlice', 'onslice-dataflow.md')
]

IMAGES = [
    'ert-overview.png',
    'onlot-dataflow.png',
    'onprod-dataflow.png',
    'onscribe-dataflow.png',
    'onslice-dataflow.png'
]

prs = Presentation()
# Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Exensio Reference Tables (ERT) — Full-res Review"
slide.placeholders[1].text = "Images embedded at full resolution; speaker notes populated from docs/*.md"

# helper funcs
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

def md_to_text(path):
    if not os.path.exists(path):
        return ''
    with open(path,'r',encoding='utf-8') as f:
        return f.read()

def md_to_bullets(path, max_bullets=8):
    if not os.path.exists(path):
        return []
    bullets = []
    with open(path,'r',encoding='utf-8') as f:
        for l in f:
            l = l.strip()
            if not l:
                continue
            if l.startswith('#'):
                continue
            if l.startswith('-') or l.startswith('*'):
                bullets.append(l.lstrip('-* ').strip())
            elif len(bullets) < max_bullets and len(l) < 120:
                bullets.append(l)
            if len(bullets) >= max_bullets:
                break
    return bullets

def add_bullet_slide(title_text, bullets, notes_text=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i>0 else tf.paragraphs[0]
        p.text = b
        p.level = 0
        p.font.size = Pt(18)
    if notes_text:
        slide.notes_slide.notes_text_frame.text = notes_text
    return slide

# create content slides with notes
for title, md in FILES:
    mdpath = os.path.join(DOCS, md)
    bullets = md_to_bullets(mdpath)
    notes = md_to_text(mdpath)
    if not bullets:
        bullets = ['See document: {}'.format(md)]
    add_bullet_slide(title, bullets[:8], notes)

# embed full-resolution images with notes describing source file
for img in IMAGES:
    imgpath = os.path.join(DOCS, img)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = os.path.splitext(img)[0].replace('-', ' ').upper()
    left = Inches(0.5)
    top = Inches(1.2)
    if os.path.exists(imgpath):
        try:
            # add_picture will embed the image at its native resolution if width/height not specified
            slide.shapes.add_picture(imgpath, left, top)
        except Exception as e:
            tx = slide.shapes.add_textbox(left, top, Inches(9), Inches(1))
            tx.text = f"[Could not render image: {img}]"
        # add notes with markdown filename
        md_equiv = os.path.splitext(img)[0].replace('-', '_') + '.md'
        mdpath = os.path.join(DOCS, md_equiv)
        if os.path.exists(mdpath):
            slide.notes_slide.notes_text_frame.text = md_to_text(mdpath)
        else:
            slide.notes_slide.notes_text_frame.text = f"Source image: {img}"
    else:
        tx = slide.shapes.add_textbox(left, top, Inches(9), Inches(1))
        tx.text = f"[Image not found: {img}]"

prs.save(OUT)
print('Wrote', OUT)
