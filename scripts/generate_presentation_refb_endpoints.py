#!/usr/bin/env python3
"""
Generate a focused presentation on REFDB endpoints and their data sources.
Endpoints covered: on_lot, on_prod, on_scribe, on_slice, on_wmap (+ optional pp_lotprod).
Slides emphasize: external data sources, key transformations, and final REFDB tables.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / 'docs'
OUT = DOCS / 'refb-endpoints-data-sources.pptx'

# Images to include (prefer PNG when available, else SVG is skipped by python-pptx)
IMAGES = {
    'OnLot': 'onlot-dataflow.png',
    'OnProd': 'onprod-dataflow.png',
    'OnScribe': 'onscribe-dataflow.png',
    'OnSlice': 'onslice-dataflow.png',
    'OnWmap (Config)': 'onwmap-level2-config.png',
    'OnWmap (Product)': 'onwmap-level2-product.png',
}

# Quick per-endpoint bullets (data sources and final tables)
BULLETS = {
    'OnLot': [
        'Sources: LotG (native + WS), LTM WS, Data Warehouse, MES (Torrent/Genesis)',
        'Config: ON_FAB_CONF, ON_SITE_CONF, ERT_CONF; trimming/heuristics for JND/Bucheon',
        'Final tables: ON_LOT, ON_PROD (persisted when enriched)',
    ],
    'OnProd': [
        'Sources (triggered by OnLot): LotG, Data Warehouse PLM/MfgArea, MES (Torrent/Genesis)',
        'Priority: MES overrides DW overrides LotG; alternate product fallback',
        'Final table: ON_PROD',
    ],
    'OnScribe': [
        'Sources: VID↔SCRIBE services (per fab), OnLot cache, AttributeUtils patterns',
        'Config: OnFabConf/ErtConf URLs, result types, onScribeWaferIdEqualsScribeId, waferIdCreationPattern',
        'Final table: ON_SCRIBE',
    ],
    'OnSlice': [
        'Sources: Direct repository (no remote), admin writes via DTO create/update',
        'Fields: slice, puckId, runId, globalWaferId, fabWaferId, source lots, part/lottype',
        'Final table: ON_SLICE',
    ],
    'OnWmap (Config)': [
        'Sources: Matchup service (MatchupLoader), WMC service via Caller',
        'Config resolution: OnFabConf or ErtConf; fetch by configuration ID',
        'Final table: ON_WMAP',
    ],
    'OnWmap (Product)': [
        'Sources: WMC by device; primary via serviceKey, secondary fallback by product',
        'Config resolution: OnFabConf/ErtConf; cached in repo if found',
        'Final table: ON_WMAP',
    ],
}

TITLE = 'Exensio REFDB endpoints — data sources to final tables'
SUB = 'Focus: ON_LOT, ON_PROD, ON_SCRIBE, ON_SLICE, ON_WMAP (+ PP_LOTPROD context)'

# Closing mapping: endpoint → final tables → key columns → primary sources
MAPPING = [
    (
        'OnLot',
        'ON_LOT; ON_PROD',
        'lot, mfgLot, product, fab, sourceLot, lotType, maskSet, process, technology, PTI, family',
        'LotG (native+WS), LTM WS (lotType), Data Warehouse (PLM/MfgArea), MES (Torrent/Genesis)'
    ),
    (
        'OnProd',
        'ON_PROD',
        'product, productVersion, family, process, technology, maskSet, fab',
        'MES (Torrent/Genesis), Data Warehouse PLM/MfgArea, LotG'
    ),
    (
        'OnScribe',
        'ON_SCRIBE',
        'lot, waferNum, waferId, scribeId, insertTime, status',
        'VID↔SCRIBE services (fab-configured); calculated fallback via AttributeUtils; OnLot cache for sourceLot context'
    ),
    (
        'OnSlice',
        'ON_SLICE',
        'slice, globalWaferId, puckId, runId, sliceSourceLot, startLot, fabWaferId, fabSourceLot, slicePartname, sliceLottype, sliceSupplierId, puckHeight, sliceOrder, sliceStartTime',
        'Primary writes via admin DTO/API; upstream ingestion uses BIWMES+eCofA+TORRENT to populate/maintain rows'
    ),
    (
        'OnWmap',
        'ON_WMAP',
        'idWaferMapConfiguration, product/device mapping, metadata per WMC',
        'Matchup service (by lot/scribe) and WMC service (by config/product) via Caller'
    ),
    (
        'PP_LOTPROD (context)',
        'PP_LOTPROD',
        'lot, product, fab (frontend provenance)',
        'Internal PP_LOTPROD DB exposed via /api/pplotprod/bylotid; consumed by ingestion scripts'
    ),
]


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = TITLE
    slide.placeholders[1].text = SUB


def add_endpoint_slide(prs, name, bullets, image_name):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    title = slide.shapes.title
    title.text = name

    # Add bullets on left
    left = Inches(0.6)
    top = Inches(1.5)
    width = Inches(5.7)
    height = Inches(5.0)
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = 'Data sources and sink:'
    p.font.size = Pt(18)
    p.font.bold = True

    for b in bullets:
        para = tf.add_paragraph()
        para.text = f"• {b}"
        para.level = 1
        para.font.size = Pt(16)

    # Add image on right if present
    img_path = DOCS / image_name if image_name else None
    if img_path and img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(6.5), Inches(1.2), height=Inches(4.8))
    else:
        # If missing image, add a placeholder note
        note = slide.shapes.add_textbox(Inches(6.3), Inches(1.2), Inches(4.5), Inches(1.0))
        nft = note.text_frame
        nft.text = f"[Diagram not found: {image_name}]"
        nft.paragraphs[0].font.color.rgb = RGBColor(200, 0, 0)


def build():
    prs = Presentation()
    add_title_slide(prs)

    # Endpoint slides
    for name, img in IMAGES.items():
        add_endpoint_slide(prs, name, BULLETS[name], img)

    # Optional PP_LOTPROD context slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = 'PP_LOTPROD (context)'
    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(10), Inches(5))
    tf = body.text_frame
    tf.word_wrap = True
    tf.text = 'Used by ingestion scripts to resolve source lot and fab via /api/pplotprod/bylotid.'
    p = tf.add_paragraph(); p.text = 'Source: internal PP_LOTPROD table via Exensio WS endpoint'; p.level = 1
    p = tf.add_paragraph(); p.text = 'Consumers: getCamstarWafer2AssemblyGenealogy.pl, getSnowflakeE142ModuleTrace.pl'; p.level = 1
    p = tf.add_paragraph(); p.text = 'Outputs: lot→product/fab context feeding ON_LOT/ON_PROD usage'; p.level = 1

    # Closing mapping slide (Endpoint → REFDB mapping)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = 'Endpoint → REFDB tables, columns, and sources'
    rows = 1 + len(MAPPING)
    cols = 4
    left, top, width, height = Inches(0.3), Inches(1.3), Inches(12.9), Inches(5.2)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    # Set headers
    headers = ['Endpoint', 'Final REFDB table(s)', 'Key columns populated', 'Primary sources']
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
    # Fill rows
    for i, (ep, tabs, cols_txt, srcs) in enumerate(MAPPING, start=1):
        table.cell(i, 0).text = ep
        table.cell(i, 1).text = tabs
        table.cell(i, 2).text = cols_txt
        table.cell(i, 3).text = srcs
        for j in range(4):
            for p in table.cell(i, j).text_frame.paragraphs:
                p.font.size = Pt(12)

    prs.save(str(OUT))
    return OUT


if __name__ == '__main__':
    out = build()
    print(f"Wrote {out}")
